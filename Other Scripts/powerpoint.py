# script to put figures in a powerpoint presentation
import os
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

typeExperiment = 'frequency' # frequency or time
imageDir = 'W:\\PhD\\PatternGlareData\\Results\\'+ typeExperiment + '\\figures'

# get the list of images
imageList = os.listdir(imageDir)
imageList = [x for x in imageList if x.endswith(".png")]


# split list into two based on if the image dir contains "0.5-3s"
imageList_3_3_99 = []
imageList_0_5_3s = []
imageList_other = []
for image in imageList:
    if "0.5-3" in image:
        imageList_0_5_3s.append(image)
    elif "2.8-3.99" in image or "3-3.99" in image:
        imageList_3_3_99.append(image)
    else:
        imageList_other.append(image)


imageList_0_5_3s_onsets = []
imageList_0_5_3s_partitions = []
imageList_0_5_3s_onsetsVS = []

imageList_3_3_99_onsets = []
imageList_3_3_99_partitions = []
imageList_3_3_99_onsetsVS = []

imageList_other_onsets = []
imageList_other_partitions = []
imageList_other_onsetsVS = []

for image in imageList_0_5_3s:
    if ("Time" in image or "by" in image) and "Onsets" in image:
        imageList_0_5_3s_onsetsVS.append(image)
    elif "Partitions" in image:
        imageList_0_5_3s_partitions.append(image)
    else:
        imageList_0_5_3s_onsets.append(image)
        
for image in imageList_3_3_99:
    if ("Time" in image or "by" in image) and "Onsets" in image:
        imageList_3_3_99_onsetsVS.append(image)
    elif "Partitions" in image:
        imageList_3_3_99_partitions.append(image)
    else:
        imageList_3_3_99_onsets.append(image)

for image in imageList_other:
    if ("Time" in image or "by" in image) and "Onsets" in image:
        imageList_other_onsetsVS.append(image)
    elif "Partitions" in image:
        imageList_other_partitions.append(image)
    else:
        imageList_other_onsets.append(image)

# sort the lists
imageList_0_5_3s_onsets.sort()
imageList_0_5_3s_partitions.sort()
imageList_0_5_3s_onsetsVS.sort()
imageList_other_onsets.sort()
imageList_other_partitions.sort()
imageList_other_onsetsVS.sort()


# create a new presentation
prs = Presentation()
# creat title slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Pattern Glare"
subtitle.text = "Results: " + typeExperiment

# creat title slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
title.text = "DC Shift Onsets"

# add the images to center of the slide but keep aspect ratio and fit to slide
for image in imageList_0_5_3s_onsets:
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    image_path = imageDir + '\\' + image
    img = Image.open(image_path)
    img_width, img_height = img.size
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    aspect_ratio = img_width / img_height
    max_height = slide_height
    new_width = int(max_height * aspect_ratio)
    left = (slide_width - new_width)/2
    top = Inches(0)
    pic = slide.shapes.add_picture(image_path, left, top, width=new_width, height=max_height)

    


# creat title slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
title.text = "DC Shift Partitions"

# add the images to center of the slide but keep aspect ratio and fit to slide
for image in imageList_0_5_3s_partitions:
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    image_path = imageDir + '\\' + image
    img = Image.open(image_path)
    img_width, img_height = img.size
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    aspect_ratio = img_width / img_height
    max_height = slide_height
    new_width = int(max_height * aspect_ratio)
    left = (slide_width - new_width)/2
    top = Inches(0)
    pic = slide.shapes.add_picture(image_path, left, top, width=new_width, height=max_height)


# creat title slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
title.text = "DC Shift Onsets 2,3 vs 4,5 vs 6,7"

# add the images to center of the slide but keep aspect ratio and fit to slide
for image in imageList_0_5_3s_onsetsVS:
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    image_path = imageDir + '\\' + image
    img = Image.open(image_path)
    img_width, img_height = img.size
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    aspect_ratio = img_width / img_height
    max_height = slide_height
    new_width = int(max_height * aspect_ratio)
    left = (slide_width - new_width)/2
    top = Inches(0)
    pic = slide.shapes.add_picture(image_path, left, top, width=new_width, height=max_height)



# creat title slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
title.text = "Offset Period Onsets"


# add the images to center of the slide but keep aspect ratio and fit to slide
for image in imageList_3_3_99_onsets:
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    image_path = imageDir + '\\' + image
    img = Image.open(image_path)
    img_width, img_height = img.size
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    aspect_ratio = img_width / img_height
    max_height = slide_height
    new_width = int(max_height * aspect_ratio)
    left = (slide_width - new_width)/2
    top = Inches(0)
    pic = slide.shapes.add_picture(image_path, left, top, width=new_width, height=max_height)

# creat title slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
title.text = "Offset Period Partitions"

# add the images to center of the slide but keep aspect ratio and fit to slide
for image in imageList_3_3_99_partitions:
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    image_path = imageDir + '\\' + image
    img = Image.open(image_path)
    img_width, img_height = img.size
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    aspect_ratio = img_width / img_height
    max_height = slide_height
    new_width = int(max_height * aspect_ratio)
    left = (slide_width - new_width)/2
    top = Inches(0)
    pic = slide.shapes.add_picture(image_path, left, top, width=new_width, height=max_height)

# creat title slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
title.text = "Offset Period Onsets 2,3 vs 4,5 vs 6,7"

# add the images to center of the slide but keep aspect ratio and fit to slide
for image in imageList_3_3_99_onsetsVS:
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    image_path = imageDir + '\\' + image
    img = Image.open(image_path)
    img_width, img_height = img.size
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    aspect_ratio = img_width / img_height
    max_height = slide_height
    new_width = int(max_height * aspect_ratio)
    left = (slide_width - new_width)/2
    top = Inches(0)
    pic = slide.shapes.add_picture(image_path, left, top, width=new_width, height=max_height)


# creat title slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
title.text = "Other Effects Onsets"

# add the images to center of the slide but keep aspect ratio and fit to slide
for image in imageList_other_onsets:
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    image_path = imageDir + '\\' + image
    img = Image.open(image_path)
    img_width, img_height = img.size
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    aspect_ratio = img_width / img_height
    max_height = slide_height
    new_width = int(max_height * aspect_ratio)
    left = (slide_width - new_width)/2
    top = Inches(0)
    pic = slide.shapes.add_picture(image_path, left, top, width=new_width, height=max_height)
    
# creat title slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
title.text = "Other Effects Partitions"

# add the images to center of the slide but keep aspect ratio and fit to slide
for image in imageList_other_partitions:
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    image_path = imageDir + '\\' + image
    img = Image.open(image_path)
    img_width, img_height = img.size
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    aspect_ratio = img_width / img_height
    max_height = slide_height
    new_width = int(max_height * aspect_ratio)
    left = (slide_width - new_width)/2
    top = Inches(0)
    pic = slide.shapes.add_picture(image_path, left, top, width=new_width, height=max_height)
    
# creat title slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
title.text = "Other Effects Onsets 2,3 vs 4,5 vs 6,7"

# add the images to center of the slide but keep aspect ratio and fit to slide
for image in imageList_other_onsetsVS:
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    image_path = imageDir + '\\' + image
    img = Image.open(image_path)
    img_width, img_height = img.size
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    aspect_ratio = img_width / img_height
    max_height = slide_height
    new_width = int(max_height * aspect_ratio)
    left = (slide_width - new_width)/2
    top = Inches(0)
    pic = slide.shapes.add_picture(image_path, left, top, width=new_width, height=max_height)

  
# save the presentation
prs.save('W:\\PhD\\PatternGlareData\\Results\\'+ typeExperiment + '\\figures\\figures.pptx')

